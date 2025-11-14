from flask import Flask, render_template_string, request, redirect, url_for, flash, send_file
import sqlite3
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io
import pandas as pd
import matplotlib.pyplot as plt
from werkzeug.utils import secure_filename
import os
import numpy as np

app = Flask(__name__)
app.secret_key = 'pitchforge_mvp_key'
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Init DB
def init_db():
    conn = sqlite3.connect('pitchforge.db')
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS pitches
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  email TEXT NOT NULL,
                  idea_summary TEXT NOT NULL,
                  target_audience TEXT,
                  funding_ask REAL,
                  timeline_months INTEGER,
                  submitted_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)''')
    conn.commit()
    conn.close()

init_db()

# Embedded form HTML
FORM_HTML = '''
<!DOCTYPE html>
<html lang="en">
<head>
    <title>PitchForge MVP - Forge Your Pitch</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet">
    <style>
        body { background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%); min-height: 100vh; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; }
        .form-card { background: white; border-radius: 15px; box-shadow: 0 10px 30px rgba(0,0,0,0.1); padding: 30px; max-width: 600px; margin: 50px auto; }
        .btn-forge { background: linear-gradient(45deg, #007bff, #0056b3); border: none; font-weight: bold; }
        .flash { color: red; }
    </style>
</head>
<body>
    <div class="container">
        <div class="form-card">
            <h1 class="text-center mb-4 text-primary">PitchForge MVP: Forge your business idea investment pitch</h1>
            {% with messages = get_flashed_messages() %}
                {% if messages %}
                    <ul class="flash">
                        {% for message in messages %}
                            <li>{{ message }}</li>
                        {% endfor %}
                    </ul>
                {% endif %}
            {% endwith %}
            <form method="POST">
                <div class="mb-3">
                    <label class="form-label fw-bold">Email</label>
                    <input type="email" class="form-control" name="email" required>
                </div>
                <div class="mb-3">
                    <label class="form-label fw-bold">Idea Summary (keep it punchy!)</label>
                    <textarea class="form-control" name="idea_summary" rows="3" required></textarea>
                </div>
                <div class="mb-3">
                    <label class="form-label fw-bold">Target Audience</label>
                    <input type="text" class="form-control" name="target_audience" placeholder="e.g., VCs in startups">
                </div>
                <div class="mb-3">
                    <label class="form-label fw-bold">Funding Ask ($)</label>
                    <input type="number" class="form-control" name="funding_ask" step="0.01">
                </div>
                <div class="mb-3">
                    <label class="form-label fw-bold">Timeline (months)</label>
                    <input type="number" class="form-control" name="timeline_months" min="1">
                </div>
                <button type="submit" class="btn btn-primary btn-forge w-100 py-3">Forge It! 🚀</button>
            </form>
            <p class="text-center mt-4 small text-muted">Investment-ready deck generated instantly—test your pitch today.</p>
        </div>
    </div>
    <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/js/bootstrap.bundle.min.js"></script>
</body>
</html>
'''

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        email = request.form['email']
        idea_summary = request.form['idea_summary']
        target_audience = request.form.get('target_audience', '')
        funding_ask_str = request.form.get('funding_ask', '')
        timeline_months_str = request.form.get('timeline_months', '')
        
        # Validation
        if not email or not idea_summary:
            flash('Email and idea summary are required!')
            return render_template_string(FORM_HTML)
        
        # Safe numbers
        try:
            funding_ask = float(funding_ask_str) if funding_ask_str else 0.0
            timeline_months = int(timeline_months_str) if timeline_months_str else 0
        except ValueError:
            flash('Funding ask and timeline must be valid numbers!')
            return render_template_string(FORM_HTML)
        
        # Save to DB
        conn = sqlite3.connect('pitchforge.db')
        c = conn.cursor()
        c.execute('''INSERT INTO pitches (email, idea_summary, target_audience, funding_ask, timeline_months)
                     VALUES (?, ?, ?, ?, ?)''',
                  (email, idea_summary, target_audience, funding_ask, timeline_months))
        pitch_id = c.lastrowid
        conn.commit()
        conn.close()
        
        # Build & stream deck
        try:
            buffer = build_pitch_deck_buffer(pitch_id, idea_summary, target_audience, funding_ask, timeline_months)
            return send_file(buffer, as_attachment=True, download_name=f'PitchForge_Deck_{pitch_id}.pptx',
                             mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        except Exception as e:
            flash(f'Deck gen error: {str(e)}—pitch saved.')
            return redirect(url_for('success'))
    
    return render_template_string(FORM_HTML)

def build_pitch_deck_buffer(pitch_id, summary, audience, ask, timeline):
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "PitchForge Investment Deck"
    subtitle = slide.placeholders[1]
    subtitle.text = f"ID: {pitch_id} | {datetime.now().strftime('%Y-%m-%d')}"
    
    # Slide 2: Idea
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Business Idea"
    content = slide.placeholders[1]
    content.text = summary
    
    # Slide 3: Market
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Target Market"
    content = slide.placeholders[1]
    content.text = f"Audience: {audience}"
    
    # Slide 4: Ask with chart
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Investment Ask"
    content = slide.placeholders[1]
    content.text = f"${ask:,.2f} | Timeline: {timeline} months"
    # Simple chart
    chart_buffer = io.BytesIO()
    plt.figure(figsize=(6, 4))
    categories = ['Funding Ask ($)', 'Timeline (Months)']
    values = [ask, timeline]
    plt.bar(categories, values, color=['blue', 'green'])
    plt.title('Investment Snapshot')
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig(chart_buffer, format='png', bbox_inches='tight')
    chart_buffer.seek(0)
    plt.close()
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(chart_buffer, left, top, width=Inches(5))
    
    # Buffer save
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

@app.route('/success')
def success():
    return '''
    <html>
        <head><title>PitchForge MVP</title></head>
        <body>
            <h1>Success! Your pitch is forged.</h1>
            <p>Deck downloaded—ready for investors. <a href="/">Submit another?</a></p>
        </body>
    </html>
    '''

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
